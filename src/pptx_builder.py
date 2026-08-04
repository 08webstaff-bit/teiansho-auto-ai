"""STEP 4: python-pptx で提案スライド（.pptx）を生成する。

構成:
  1. 表紙
  2. 提案コンセプト
  3. 解決策のイメージ
  4. 参考類似事例（1〜2ページ・写真＋URL＋選定理由）
  5. お見積内容（表＋合計）
  6. 裏表紙

デザイン: 白ベース＋コーポレートカラー #0070E0、游ゴシック。
"""

import io

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt

from .scrape import fetch_case_images

CORPORATE = RGBColor(0x00, 0x70, 0xE0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x88, 0x88, 0x88)
LIGHT = RGBColor(0xF0, 0xF4, 0xFA)
FONT = "游ゴシック"

# 16:9 スライドサイズ
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

COMPANY_NAME = "株式会社丸八テント商会"
COMPANY_URL = "https://08tent.co.jp/"
COMPANY_TEL = "お問い合わせは担当営業までご連絡ください"


def _set_font(run, size=18, bold=False, color=DARK):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # 日本語フォントを East Asian にも適用
    rpr = run._r.get_or_add_rPr()
    from pptx.oxml.ns import qn

    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = rpr.makeelement(qn("a:ea"), {})
        rpr.append(ea)
    ea.set("typeface", FONT)


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 完全空白レイアウト


def _add_textbox(slide, left, top, width, height, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return box, tf


def _add_paragraph(tf, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, first=False, space_after=6):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    _set_font(run, size=size, bold=bold, color=color)
    return p


def _add_rect(slide, left, top, width, height, color):
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _fmt_yen(value):
    try:
        return f"¥{int(value):,}"
    except (TypeError, ValueError):
        return "―"


# ------------------------------------------------------------------ 各スライド


def _slide_cover(prs, quote, text):
    slide = _blank_slide(prs)
    _add_rect(slide, 0, 0, SLIDE_W, Emu(120000), CORPORATE)  # 上部ライン
    _add_rect(slide, 0, SLIDE_H - Emu(120000), SLIDE_W, Emu(120000), CORPORATE)  # 下部ライン

    customer = quote.get("customer_name") or "お客様"
    _, tf = _add_textbox(slide, Emu(1000000), Emu(2200000), Emu(10192000), Emu(2000000), anchor=MSO_ANCHOR.MIDDLE)
    _add_paragraph(tf, f"{customer} 御中", size=28, bold=True, color=DARK, first=True, space_after=18)
    _add_paragraph(tf, "ご 提 案 書", size=44, bold=True, color=CORPORATE, space_after=18)
    title = text.get("title_suggestion") or quote.get("project_name") or ""
    if title:
        _add_paragraph(tf, title, size=22, color=GRAY)

    _, tf2 = _add_textbox(slide, Emu(1000000), SLIDE_H - Emu(1100000), Emu(10192000), Emu(600000))
    _add_paragraph(tf2, f"提案：{COMPANY_NAME}", size=18, bold=True, color=DARK, first=True, align=PP_ALIGN.RIGHT)


def _header(slide, title):
    _add_rect(slide, 0, 0, Emu(160000), SLIDE_H, CORPORATE)  # 左サイドバー
    _, tf = _add_textbox(slide, Emu(600000), Emu(360000), Emu(11000000), Emu(700000))
    _add_paragraph(tf, title, size=28, bold=True, color=CORPORATE, first=True)
    _add_rect(slide, Emu(600000), Emu(1050000), Emu(2400000), Emu(45000), CORPORATE)


def _slide_concept(prs, text):
    slide = _blank_slide(prs)
    _header(slide, "提案コンセプト")
    _, tf = _add_textbox(slide, Emu(700000), Emu(1500000), Emu(10800000), Emu(4400000))
    _add_paragraph(tf, text.get("concept", ""), size=20, color=DARK, first=True, space_after=12)


def _slide_solution(prs, text):
    slide = _blank_slide(prs)
    _header(slide, "解決策のイメージ")
    top = Emu(1650000)
    for i, item in enumerate(text.get("solution_images", [])):
        card = _add_rect(slide, Emu(700000), top, Emu(10800000), Emu(1500000), LIGHT)
        _, tf = _add_textbox(slide, Emu(1000000), top + Emu(300000), Emu(10200000), Emu(900000), anchor=MSO_ANCHOR.MIDDLE)
        _add_paragraph(tf, f"施工イメージ案 {i + 1}", size=15, bold=True, color=CORPORATE, first=True, space_after=6)
        _add_paragraph(tf, item, size=19, color=DARK)
        top += Emu(1800000)


def _slide_case(prs, resolved_case):
    slide = _blank_slide(prs)
    _header(slide, "参考類似事例")

    # 写真（メイン）を取得して掲載。失敗時はプレースホルダー枠。
    images = []
    if "/works/" in resolved_case.get("url", ""):
        try:
            images = fetch_case_images(resolved_case["url"], max_images=2)
        except Exception:
            images = []

    img_left, img_top, img_w = Emu(700000), Emu(1500000), Emu(5400000)
    if images:
        try:
            slide.shapes.add_picture(images[0], img_left, img_top, width=img_w)
        except Exception:
            _placeholder(slide, img_left, img_top, img_w, Emu(3600000))
    else:
        _placeholder(slide, img_left, img_top, img_w, Emu(3600000))

    # 右側テキスト
    _, tf = _add_textbox(slide, Emu(6400000), Emu(1500000), Emu(5100000), Emu(4200000))
    _add_paragraph(tf, resolved_case.get("title", ""), size=22, bold=True, color=DARK, first=True, space_after=16)
    if resolved_case.get("category_name") and resolved_case.get("is_individual"):
        _add_paragraph(tf, f"カテゴリー: {resolved_case['category_name']}", size=13, color=GRAY, space_after=20)
    # 選定理由は社内向け情報のため、顧客向けスライドには掲載しない
    _add_paragraph(tf, "施工事例ページ:", size=13, bold=True, color=GRAY, space_after=2)
    url_p = tf.add_paragraph()
    run = url_p.add_run()
    run.text = resolved_case.get("url", "")
    _set_font(run, size=13, color=CORPORATE)
    _hyperlink(run, resolved_case.get("url", ""))

    # 2枚目のサブ写真があれば下に小さく
    if len(images) > 1:
        try:
            slide.shapes.add_picture(images[1], img_left, Emu(5200000), width=Emu(2600000))
        except Exception:
            pass


def _placeholder(slide, left, top, width, height):
    _add_rect(slide, left, top, width, height, LIGHT)
    _, tf = _add_textbox(slide, left, top, width, height, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _add_paragraph(tf, "（写真は施工事例ページをご覧ください）", size=14, color=GRAY, first=True, align=PP_ALIGN.CENTER)


def _hyperlink(run, url):
    if url:
        run.hyperlink.address = url


def _slide_quote(prs, quote):
    slide = _blank_slide(prs)
    _header(slide, "お見積内容")

    items = quote.get("items", []) or []
    rows = len(items) + 2  # ヘッダー＋明細＋合計
    cols = 5
    left, top = Emu(600000), Emu(1450000)
    width, height = Emu(11000000), Emu(min(4600000, 500000 + rows * 360000))
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    table.columns[0].width = Emu(4200000)
    table.columns[1].width = Emu(3200000)
    table.columns[2].width = Emu(1200000)
    table.columns[3].width = Emu(1200000)
    table.columns[4].width = Emu(1200000)

    headers = ["品名", "仕様", "数量", "単価", "金額"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CORPORATE
        _cell_text(cell, h, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for r, item in enumerate(items, start=1):
        _cell_text(table.cell(r, 0), item.get("name", ""), size=13)
        _cell_text(table.cell(r, 1), item.get("spec", ""), size=12)
        _cell_text(table.cell(r, 2), str(item.get("quantity", "")), size=13, align=PP_ALIGN.CENTER)
        _cell_text(table.cell(r, 3), _fmt_yen(item.get("unit_price")), size=13, align=PP_ALIGN.RIGHT)
        _cell_text(table.cell(r, 4), _fmt_yen(item.get("amount")), size=13, align=PP_ALIGN.RIGHT)
        for c in range(cols):
            table.cell(r, c).fill.solid()
            table.cell(r, c).fill.fore_color.rgb = WHITE if r % 2 else LIGHT

    total_row = rows - 1
    table.cell(total_row, 0).merge(table.cell(total_row, 3))
    _cell_text(table.cell(total_row, 0), "合計金額（税抜）", size=15, bold=True, align=PP_ALIGN.RIGHT)
    _cell_text(table.cell(total_row, 4), _fmt_yen(quote.get("total_amount")), size=15, bold=True, color=CORPORATE, align=PP_ALIGN.RIGHT)
    for c in range(cols):
        table.cell(total_row, c).fill.solid()
        table.cell(total_row, c).fill.fore_color.rgb = LIGHT


def _cell_text(cell, text, size=13, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Emu(90000)
    cell.margin_right = Emu(90000)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_font(run, size=size, bold=bold, color=color)


def _slide_back(prs):
    slide = _blank_slide(prs)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, CORPORATE)
    _, tf = _add_textbox(slide, Emu(1000000), Emu(2400000), Emu(10192000), Emu(2200000), anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    _add_paragraph(tf, COMPANY_NAME, size=34, bold=True, color=WHITE, first=True, align=PP_ALIGN.CENTER, space_after=24)
    _add_paragraph(tf, COMPANY_URL, size=20, color=WHITE, align=PP_ALIGN.CENTER, space_after=12)
    _add_paragraph(tf, COMPANY_TEL, size=16, color=WHITE, align=PP_ALIGN.CENTER)


# ------------------------------------------------------------------ エントリ


def build_proposal_pptx(quote, resolved_cases, proposal_text, output=None):
    """提案書 pptx を生成する。

    output が None の場合は bytes を返す。パス文字列なら保存してパスを返す。
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide_cover(prs, quote, proposal_text)
    _slide_concept(prs, proposal_text)
    _slide_solution(prs, proposal_text)
    for case in (resolved_cases or [])[:2]:
        _slide_case(prs, case)
    _slide_quote(prs, quote)
    _slide_back(prs)

    if output is None:
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()
    prs.save(output)
    return output
