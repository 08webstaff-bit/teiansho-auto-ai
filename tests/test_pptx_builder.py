"""pptx_builder の生成ロジックの単体テスト（画像取得はモック）。"""

import io

import pytest
from pptx import Presentation

import src.pptx_builder as pb

QUOTE = {
    "customer_name": "株式会社テスト工業",
    "project_name": "荷捌き場テント新設工事",
    "industry_type": "工場",
    "total_amount": 1650000,
    "items": [
        {"name": "荷捌き場テント", "spec": "片持ち式 W6000", "quantity": "1式", "unit_price": 1500000, "amount": 1500000},
        {"name": "施工費", "spec": "", "quantity": "1式", "unit_price": 150000, "amount": 150000},
    ],
}

RESOLVED = [
    {"key": "nisabaki_tent", "category_name": "荷捌き場テント一覧", "url": "https://08tent.co.jp/works/83528/", "title": "大型トラック対応 片持ちテント屋根", "thumbnail": None, "reason": "仕様一致", "is_individual": True, "resolved": True, "category_reason": "荷捌き用途"},
    {"key": "jabara_tent", "category_name": "工場間通路テント", "url": "https://08tent.co.jp/works/58612/", "title": "工場間通路テント（常設）", "thumbnail": None, "reason": "", "is_individual": False, "resolved": True, "category_reason": "常設ニーズ"},
]

TEXT = {
    "title_suggestion": "荷捌き場テント新設のご提案",
    "concept": "この度はお見積のご依頼をありがとうございます。" * 4,
    "solution_images": ["片持ち式で車両動線を確保", "照明付きで夜間作業対応"],
}


@pytest.fixture(autouse=True)
def _no_network(monkeypatch):
    """画像取得はネットワークを使わずスタブ化。"""
    monkeypatch.setattr(pb, "fetch_case_images", lambda url, max_images=2: [])


def test_build_returns_bytes():
    data = pb.build_proposal_pptx(QUOTE, RESOLVED, TEXT)
    assert isinstance(data, bytes) and len(data) > 5000


def test_slide_count_with_two_cases():
    data = pb.build_proposal_pptx(QUOTE, RESOLVED, TEXT)
    prs = Presentation(io.BytesIO(data))
    # 表紙+コンセプト+解決策+事例2+見積+裏表紙 = 7
    assert len(prs.slides) == 7


def test_slide_count_with_one_case():
    data = pb.build_proposal_pptx(QUOTE, RESOLVED[:1], TEXT)
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) == 6


def test_cover_contains_customer_and_title():
    data = pb.build_proposal_pptx(QUOTE, RESOLVED, TEXT)
    prs = Presentation(io.BytesIO(data))
    texts = _all_text(prs.slides[0])
    assert "株式会社テスト工業 御中" in texts
    assert "ご 提 案 書" in texts


def test_quote_slide_has_total_amount():
    data = pb.build_proposal_pptx(QUOTE, RESOLVED, TEXT)
    prs = Presentation(io.BytesIO(data))
    all_text = " ".join(_all_text(s) for s in prs.slides)
    assert "¥1,650,000" in all_text
    assert "荷捌き場テント" in all_text


def test_case_slide_has_url():
    data = pb.build_proposal_pptx(QUOTE, RESOLVED, TEXT)
    prs = Presentation(io.BytesIO(data))
    all_text = " ".join(_all_text(s) for s in prs.slides)
    assert "https://08tent.co.jp/works/83528/" in all_text


def test_handles_empty_quote_gracefully():
    empty = {"customer_name": "", "project_name": "", "industry_type": "その他", "total_amount": None, "items": []}
    data = pb.build_proposal_pptx(empty, [], TEXT)
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) == 5  # 事例スライドなし


def _all_text(slide):
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    parts.append(cell.text)
    return " ".join(parts)
